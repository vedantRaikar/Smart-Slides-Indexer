"""Worker pipeline for PPTX indexing.

Implements a stateless worker pipeline with stages:
1. Parser -> Extract slides from PPTX
2. OCR -> Extract text from images
3. Structure -> Analyze document structure
4. Metadata (LLM) -> Extract metadata, summaries
5. Embed -> Generate embeddings
6. Graph -> Build slide relationship graph

Features:
- Idempotent processing with job IDs
- Retry with exponential backoff
- Parallel processing
- Intermediate artifact persistence
"""

import hashlib
import json
import logging
import time
import uuid
from dataclasses import dataclass, field
from enum import Enum
from pathlib import Path
from typing import Any, Callable, Dict, List, Optional

from pptx_indexer.config import get_config
from pptx_indexer.llm_adapter import LLMAdapter, create_llm_adapter
from pptx_indexer.vector_store import BaseVectorStore, create_vector_store

logger = logging.getLogger(__name__)


class JobStatus(str, Enum):
    """Job status enumeration."""

    PENDING = "pending"
    RUNNING = "running"
    COMPLETED = "completed"
    FAILED = "failed"
    RETRYING = "retrying"


@dataclass
class JobState:
    """Job state for tracking progress."""

    job_id: str
    status: JobStatus = JobStatus.PENDING
    current_stage: Optional[str] = None
    progress: float = 0.0
    error: Optional[str] = None
    stages_completed: List[str] = field(default_factory=list)
    metadata: Dict[str, Any] = field(default_factory=dict)
    created_at: float = field(default_factory=time.time)
    updated_at: float = field(default_factory=time.time)

    def to_dict(self) -> Dict:
        return {
            "job_id": self.job_id,
            "status": self.status.value,
            "current_stage": self.current_stage,
            "progress": self.progress,
            "error": self.error,
            "stages_completed": self.stages_completed,
            "metadata": self.metadata,
        }


class PipelineStage:
    """Base class for pipeline stages."""

    def __init__(self, name: str):
        self.name = name

    def process(self, context: "PipelineContext") -> "PipelineContext":
        """Process the context and return updated context."""
        raise NotImplementedError


class PipelineContext:
    """Pipeline execution context shared between stages."""

    def __init__(
        self,
        job_id: str,
        input_path: str,
        output_dir: str,
        artifacts_dir: str,
    ):
        self.job_id = job_id
        self.input_path = input_path
        self.output_dir = Path(output_dir)
        self.artifacts_dir = Path(artifacts_dir)
        self.artifacts_dir.mkdir(parents=True, exist_ok=True)

        # Data accumulated through stages
        self.slides: List[Dict] = []
        self.embeddings: Dict[str, List[float]] = {}
        self.graph: Dict = {}
        self.metadata: Dict = {}
        self.errors: List[str] = []
        self.stages_completed: List[str] = []

    def save_artifact(self, name: str, data: Any):
        """Save intermediate artifact."""
        artifact_path = self.artifacts_dir / f"{self.job_id}_{name}.json"
        if isinstance(data, (dict, list)):
            artifact_path.write_text(json.dumps(data, indent=2))
        else:
            artifact_path.write_text(str(data))
        logger.debug(f"Saved artifact: {artifact_path}")

    def load_artifact(self, name: str) -> Optional[Any]:
        """Load intermediate artifact."""
        artifact_path = self.artifacts_dir / f"{self.job_id}_{name}.json"
        if artifact_path.exists():
            return json.loads(artifact_path.read_text())
        return None

    def compute_content_hash(self) -> str:
        """Compute hash of slide content for idempotency."""
        content = json.dumps(self.slides, sort_keys=True)
        return hashlib.sha256(content.encode()).hexdigest()


class Retryable:
    """Decorator for retryable operations."""

    def __init__(self, max_attempts: int = 3, backoff: float = 1.0):
        self.max_attempts = max_attempts
        self.backoff = backoff

    def __call__(self, func: Callable) -> Callable:
        def wrapper(*args, **kwargs):
            last_exception = None
            for attempt in range(self.max_attempts):
                try:
                    return func(*args, **kwargs)
                except Exception as e:
                    last_exception = e
                    if attempt < self.max_attempts - 1:
                        wait_time = self.backoff * (2**attempt)
                        logger.warning(
                            f"{func.__name__} failed (attempt {attempt + 1}/{self.max_attempts}): {e}. "
                            f"Retrying in {wait_time}s..."
                        )
                        time.sleep(wait_time)
                    else:
                        logger.error(f"{func.__name__} failed after {self.max_attempts} attempts")
            raise last_exception

        return wrapper


class PPTXParserStage(PipelineStage):
    """Stage 1: Parse PPTX file to slides."""

    def __init__(self):
        super().__init__("parser")

    @Retryable(max_attempts=3, backoff=1.0)
    def process(self, context: PipelineContext) -> PipelineContext:
        logger.info(f"[{self.name}] Parsing {context.input_path}")

        # Check for cached parsed data
        cached = context.load_artifact("parsed_slides")
        if cached:
            context.slides = cached
            logger.info(f"[{self.name}] Loaded cached slides: {len(context.slides)}")
            return context

        # Parse PPTX
        from pptx import Presentation

        prs = Presentation(context.input_path)
        slides = []

        for i, slide in enumerate(prs.slides):
            slide_data = {
                "slide_number": i + 1,
                "slide_id": str(uuid.uuid4()),
                "title": "",
                "bullets": [],
                "images": [],
                "tables": [],
                "notes": "",
            }

            # Extract title
            if slide.shapes.title:
                slide_data["title"] = slide.shapes.title.text

            # Extract bullets from content
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            slide_data["bullets"].append(
                                {
                                    "text": para.text,
                                    "level": para.level,
                                }
                            )

            # Extract images
            for shape in slide.shapes:
                if (
                    hasattr(shape, "shape_type")
                    and str(shape.shape_type) == "MSO_SHAPE_TYPE.PICTURE"
                ):
                    slide_data["images"].append(
                        {
                            "name": shape.name,
                            "left": shape.left,
                            "top": shape.top,
                        }
                    )

            # Extract notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                slide_data["notes"] = slide.notes_slide.notes_text_frame.text

            slides.append(slide_data)

        context.slides = slides
        context.save_artifact("parsed_slides", slides)

        logger.info(f"[{self.name}] Parsed {len(slides)} slides")
        return context


class OCRStage(PipelineStage):
    """Stage 2: OCR for images."""

    def __init__(self, ocr_provider: str = "paddleocr", languages: List[str] = None):
        super().__init__("ocr")
        self.ocr_provider = ocr_provider
        self.languages = languages or ["en"]

    @Retryable(max_attempts=3, backoff=1.0)
    def process(self, context: PipelineContext) -> PipelineContext:
        logger.info(f"[{self.name}] Running OCR on {len(context.slides)} slides")

        # Check for cached OCR data
        cached = context.load_artifact("ocr_text")
        if cached:
            for i, slide in enumerate(context.slides):
                if str(i) in cached:
                    slide["ocr_text"] = cached[str(i)]
            return context

        # Initialize OCR
        if self.ocr_provider == "paddleocr":
            from paddleocr import PaddleOCR

            ocr = PaddleOCR(lang=self.languages[0], use_angle_cls=True)
        elif self.ocr_provider == "pytesseract":
            ocr = None
        else:
            logger.warning(f"Unknown OCR provider: {self.ocr_provider}")
            return context

        ocr_results = {}

        for i, slide in enumerate(context.slides):
            if not slide.get("images"):
                continue

            # Process each image
            # Note: In real implementation, images would be extracted from PPTX
            # For now, this is a placeholder
            slide["ocr_text"] = ""
            ocr_results[str(i)] = ""

        context.save_artifact("ocr_text", ocr_results)
        logger.info(f"[{self.name}] OCR completed")
        return context


class StructureAnalyzerStage(PipelineStage):
    """Stage 3: Analyze document structure."""

    def __init__(self):
        super().__init__("structure")

    def process(self, context: PipelineContext) -> PipelineContext:
        logger.info(f"[{self.name}] Analyzing structure")

        # Analyze structure from slides
        sections = []
        current_section = {
            "id": str(uuid.uuid4()),
            "title": "Introduction",
            "slide_ids": [],
        }

        for slide in context.slides:
            slide_id = slide["slide_id"]

            # Simple heuristic: new section when title changes significantly
            # In production, use LLM for section detection
            if slide["slide_number"] == 1:
                current_section["title"] = slide.get("title", "Introduction")

            current_section["slide_ids"].append(slide_id)

            # Create new section on major content boundary
            if "summary" in slide.get("title", "").lower():
                sections.append(current_section)
                current_section = {
                    "id": str(uuid.uuid4()),
                    "title": "Next Section",
                    "slide_ids": [],
                }

        if current_section["slide_ids"]:
            sections.append(current_section)

        context.metadata["sections"] = sections
        context.metadata["total_sections"] = len(sections)

        logger.info(f"[{self.name}] Found {len(sections)} sections")
        return context


class MetadataExtractionStage(PipelineStage):
    """Stage 4: Extract metadata using LLM."""

    def __init__(self, llm_adapter: LLMAdapter = None):
        super().__init__("metadata")
        self.llm = llm_adapter or create_llm_adapter()

    def process(self, context: PipelineContext) -> PipelineContext:
        logger.info(f"[{self.name}] Extracting metadata with LLM")

        # Check for cached metadata
        cached = context.load_artifact("metadata")
        if cached:
            for slide in context.slides:
                slide_id = slide["slide_id"]
                if slide_id in cached:
                    slide.update(cached[slide_id])
            return context

        metadata = {}

        # Process in batches
        batch_size = get_config().llm.batch_size

        for i in range(0, len(context.slides), batch_size):
            batch = context.slides[i : i + batch_size]

            prompts = []
            for slide in batch:
                slide_text = f"Title: {slide.get('title', '')}\n"
                slide_text += "Content: " + " | ".join(
                    b["text"] for b in slide.get("bullets", [])[:5]
                )

                prompt = f"""Extract keywords and learning objectives from this slide.
                
Slide: {slide_text}

Return JSON with:
- keywords: list of 5-10 relevant keywords
- learning_objectives: list of 3-5 learning objectives

JSON:"""
                prompts.append(prompt)

            # Batch generate
            if self.llm:
                try:
                    responses = self.llm.batch_generate(prompts)

                    for j, slide in enumerate(batch):
                        slide_id = slide["slide_id"]
                        try:
                            # Parse JSON response
                            import re

                            json_match = re.search(r"\{.*\}", responses[j].text, re.DOTALL)
                            if json_match:
                                extracted = json.loads(json_match.group())
                                slide["keywords"] = extracted.get("keywords", [])
                                slide["learning_objectives"] = extracted.get(
                                    "learning_objectives", []
                                )
                                metadata[slide_id] = extracted
                        except (json.JSONDecodeError, AttributeError) as e:
                            logger.warning(
                                f"Failed to parse LLM response for slide {slide_id}: {e}"
                            )
                except Exception as e:
                    logger.error(f"LLM batch generation failed: {e}")

        context.save_artifact("metadata", metadata)
        logger.info(f"[{self.name}] Metadata extraction completed")
        return context


class EmbeddingStage(PipelineStage):
    """Stage 5: Generate embeddings."""

    def __init__(
        self,
        embedder=None,
        batch_size: int = 32,
        cache_enabled: bool = True,
    ):
        super().__init__("embedding")
        self.embedder = embedder
        self.batch_size = batch_size
        self.cache_enabled = cache_enabled

        if not embedder:
            from sentence_transformers import SentenceTransformer

            self.embedder = SentenceTransformer("all-MiniLM-L6-v2")

    def process(self, context: PipelineContext) -> PipelineContext:
        logger.info(f"[{self.name}] Generating embeddings")

        # Check for cached embeddings
        cached = context.load_artifact("embeddings")
        if cached:
            context.embeddings = cached
            return context

        # Prepare texts for embedding
        texts = []
        for slide in context.slides:
            text = f"Title: {slide.get('title', '')}\n"
            text += " | ".join(b["text"] for b in slide.get("bullets", [])[:5])
            texts.append(text)

        # Generate embeddings in batches
        embeddings = self.embedder.encode(
            texts,
            batch_size=self.batch_size,
            show_progress_bar=True,
        )

        # Store embeddings with slide IDs
        for slide, emb in zip(context.slides, embeddings):
            context.embeddings[slide["slide_id"]] = emb.tolist()

        context.save_artifact("embeddings", context.embeddings)
        logger.info(f"[{self.name}] Generated {len(context.embeddings)} embeddings")
        return context


class GraphBuilderStage(PipelineStage):
    """Stage 6: Build slide relationship graph."""

    def __init__(self):
        super().__init__("graph")

    def process(self, context: PipelineContext) -> PipelineContext:
        logger.info(f"[{self.name}] Building slide graph")

        import numpy as np
        from sklearn.metrics.pairwise import cosine_similarity

        # Build adjacency based on content similarity
        slide_ids = list(context.embeddings.keys())
        embeddings = np.array([context.embeddings[sid] for sid in slide_ids])

        # Compute similarity matrix
        sim_matrix = cosine_similarity(embeddings)

        # Build graph
        nodes = []
        edges = []
        threshold = 0.3  # Similarity threshold

        for i, slide in enumerate(context.slides):
            nodes.append(
                {
                    "id": slide["slide_id"],
                    "slide_number": slide["slide_number"],
                    "title": slide.get("title", ""),
                }
            )

        for i in range(len(slide_ids)):
            for j in range(i + 1, len(slide_ids)):
                if sim_matrix[i][j] > threshold:
                    edges.append(
                        {
                            "source": slide_ids[i],
                            "target": slide_ids[j],
                            "weight": float(sim_matrix[i][j]),
                        }
                    )

        context.graph = {
            "version": "1.0",
            "nodes": nodes,
            "edges": edges,
            "stats": {
                "total_nodes": len(nodes),
                "total_edges": len(edges),
            },
        }

        logger.info(f"[{self.name}] Built graph with {len(nodes)} nodes, {len(edges)} edges")
        return context


class IndexerWorker:
    """Main indexing worker orchestrating the pipeline."""

    def __init__(
        self,
        llm_adapter: LLMAdapter = None,
        vector_store: BaseVectorStore = None,
        max_workers: int = None,
    ):
        config = get_config()

        self.llm = llm_adapter or create_llm_adapter()
        self.vector_store = vector_store or create_vector_store()
        self.max_workers = max_workers or config.worker.max_workers

        # Initialize pipeline stages
        self.stages: List[PipelineStage] = [
            PPTXParserStage(),
            OCRStage(
                ocr_provider=config.ocr.provider,
                languages=config.ocr.languages,
            ),
            StructureAnalyzerStage(),
            MetadataExtractionStage(self.llm),
            EmbeddingStage(),
            GraphBuilderStage(),
        ]

    def process(
        self,
        input_path: str,
        output_dir: str = "./data/output",
        job_id: Optional[str] = None,
    ) -> PipelineContext:
        """Process a PPTX file through the full pipeline."""
        # Generate job ID if not provided
        job_id = job_id or str(uuid.uuid4())

        # Check for idempotency
        config = get_config()
        artifacts_dir = config.worker.artifacts_path

        # Create context
        context = PipelineContext(
            job_id=job_id,
            input_path=input_path,
            output_dir=output_dir,
            artifacts_dir=artifacts_dir,
        )

        logger.info(f"Starting indexing job: {job_id}")

        try:
            # Run pipeline stages
            for stage in self.stages:
                if config.ocr.enabled is False and stage.name == "ocr":
                    continue

                context.current_stage = stage.name
                logger.info(f"Running stage: {stage.name}")

                stage.process(context)
                context.stages_completed.append(stage.name)

                # Update progress
                progress = len(context.stages_completed) / len(self.stages)
                logger.info(f"Progress: {progress * 100:.1f}%")

            # Final output
            self._write_output(context, output_dir)

            # Store in vector store
            self._store_vectors(context)

            logger.info(f"Job {job_id} completed successfully")
            return context

        except Exception as e:
            logger.error(f"Job {job_id} failed: {e}")
            context.errors.append(str(e))
            raise

    def _write_output(self, context: PipelineContext, output_dir: str):
        """Write final index and graph to output directory."""
        output_path = Path(output_dir)
        output_path.mkdir(parents=True, exist_ok=True)

        # Write index.json
        index_data = {
            "document_id": context.job_id,
            "document_title": Path(context.input_path).stem,
            "stats": {
                "total_slides": len(context.slides),
                "total_sections": context.metadata.get("total_sections", 0),
                "total_images": sum(len(s.get("images", [])) for s in context.slides),
            },
            "slides": {s["slide_id"]: s for s in context.slides},
            "sections": {s["id"]: s for s in context.metadata.get("sections", [])},
        }

        (output_path / "index.json").write_text(json.dumps(index_data, indent=2))

        # Write graph.json
        (output_path / "graph.json").write_text(json.dumps(context.graph, indent=2))

        logger.info(f"Output written to {output_dir}")

    def _store_vectors(self, context: PipelineContext):
        """Store embeddings in vector store."""
        ids = []
        embeddings = []
        texts = []
        metadata = []

        for slide in context.slides:
            slide_id = slide["slide_id"]
            if slide_id not in context.embeddings:
                continue

            text = f"{slide.get('title', '')} | " + " | ".join(
                b["text"] for b in slide.get("bullets", [])[:5]
            )

            ids.append(slide_id)
            embeddings.append(context.embeddings[slide_id])
            texts.append(text)
            metadata.append(
                {
                    "slide_number": slide["slide_number"],
                    "title": slide.get("title", ""),
                }
            )

        if ids:
            self.vector_store.add(ids, embeddings, texts, metadata)
            self.vector_store.persist()
            logger.info(f"Stored {len(ids)} vectors")


def run_indexing(
    input_path: str,
    output_dir: str = "./data/output",
    job_id: Optional[str] = None,
) -> Dict:
    """Convenience function to run indexing."""
    worker = IndexerWorker()
    context = worker.process(input_path, output_dir, job_id)
    return {
        "job_id": context.job_id,
        "status": "completed",
        "slides": len(context.slides),
        "sections": context.metadata.get("total_sections", 0),
        "output_dir": output_dir,
    }
