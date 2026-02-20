"""
PPT Parser - Extracts structured content from PowerPoint files.
Foundation module for the entire indexing framework.
"""

import os
from typing import List, Dict, Any, Optional, Tuple
from pathlib import Path
import logging

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE

from ..schemas.slide_node import (
    SlideNode,
    BulletPoint,
    ImageNode,
    TableNode,
)

logger = logging.getLogger(__name__)


class PPTParser:
    """
    Parses PowerPoint presentations into structured format.
    Preserves hierarchy, layout, and semantic information.
    """
    
    def __init__(self, output_dir: str = "./extracted_images"):
        self.output_dir = output_dir
        Path(output_dir).mkdir(exist_ok=True)
    
    def parse(self, pptx_path: str) -> List[SlideNode]:
        """
        Parse entire presentation into structured slides.
        
        Args:
            pptx_path: Path to PowerPoint file
        
        Returns:
            List of parsed SlideNode objects
        """
        logger.info(f"Parsing presentation: {pptx_path}")
        
        try:
            presentation = Presentation(pptx_path)
        except Exception as e:
            logger.error(f"Failed to load presentation: {e}")
            raise
        
        slides = []
        for slide_num, prs_slide in enumerate(presentation.slides, start=1):
            try:
                slide = self._parse_slide(prs_slide, slide_num, pptx_path)
                slides.append(slide)
                logger.debug(f"Parsed slide {slide_num}: {slide.title}")
            except Exception as e:
                logger.warning(f"Failed to parse slide {slide_num}: {e}")
                continue
        
        logger.info(f"Successfully parsed {len(slides)} slides")
        return slides
    
    def _parse_slide(
        self,
        prs_slide,
        slide_number: int,
        source_path: str,
    ) -> SlideNode:
        """Parse individual slide."""
        slide = SlideNode(slide_number=slide_number)
        
        # Extract layout type
        slide.layout_type = prs_slide.slide_layout.name if prs_slide.slide_layout else "default"
        
        # Process all shapes on slide
        for shape in prs_slide.shapes:
            self._process_shape(shape, slide, source_path, slide_number)
        
        # Extract speaker notes
        if prs_slide.has_notes_slide:
            notes_frame = prs_slide.notes_slide.notes_text_frame
            if notes_frame:
                slide.notes = notes_frame.text
        
        return slide
    
    def _process_shape(
        self,
        shape,
        slide: SlideNode,
        source_path: str,
        slide_number: int,
    ) -> None:
        """Process individual shape (text, image, table, etc)."""
        
        # Handle text boxes and text frames
        if shape.has_text_frame:
            self._extract_text(shape, slide)
        
        # Handle images
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            self._extract_image(shape, slide, source_path, slide_number)
        
        # Handle tables
        elif shape.has_table:
            self._extract_table(shape, slide)
    
    def _extract_text(self, shape, slide: SlideNode) -> None:
        """Extract text content preserving hierarchy."""
        text_frame = shape.text_frame
        full_text = text_frame.text.strip()
        
        if not full_text:
            return
        
        # Determine if this is title/subtitle or content
        is_title = "TITLE" in shape.name.upper()
        is_subtitle = "SUBTITLE" in shape.name.upper()
        
        if is_title and not slide.title:
            slide.title = full_text
            logger.debug(f"Extracted title: {full_text[:50]}")
            
        elif is_subtitle and not slide.subtitle:
            slide.subtitle = full_text
            logger.debug(f"Extracted subtitle: {full_text[:50]}")
            
        else:
            # Extract bullet points preserving hierarchy
            for paragraph in text_frame.paragraphs:
                text = paragraph.text.strip()
                if text:
                    level = paragraph.level
                    bullet = BulletPoint(
                        text=text,
                        level=level,
                        index=len(slide.bullets),
                    )
                    slide.bullets.append(bullet)
                    logger.debug(f"Extracted bullet (level {level}): {text[:40]}")
    
    def _extract_image(
        self,
        shape,
        slide: SlideNode,
        source_path: str,
        slide_number: int,
    ) -> None:
        """Extract and save images from slide."""
        try:
            image = shape.image
            image_bytes = image.blob
            
            # Generate image filename
            filename = f"slide_{slide_number:03d}_img_{len(slide.images)}.{image.ext}"
            image_path = os.path.join(self.output_dir, filename)
            
            # Save image
            with open(image_path, "wb") as f:
                f.write(image_bytes)
            
            # Create image node
            image_node = ImageNode(
                image_id=f"img_{slide_number}_{len(slide.images)}",
                image_path=image_path,
            )
            slide.images.append(image_node)
            logger.debug(f"Extracted image: {image_path}")
            
        except Exception as e:
            logger.warning(f"Failed to extract image: {e}")
    
    def _extract_table(self, shape, slide: SlideNode) -> None:
        """Extract table data from slide."""
        try:
            table = shape.table
            
            # Extract headers
            headers = []
            for cell in table.rows[0].cells:
                headers.append(cell.text.strip())
            
            # Extract rows
            rows = []
            for row in table.rows[1:]:
                row_data = []
                for cell in row.cells:
                    row_data.append(cell.text.strip())
                rows.append(row_data)
            
            # Create table node
            table_node = TableNode(
                table_id=f"table_{len(slide.tables)}",
                headers=headers,
                rows=rows,
            )
            slide.tables.append(table_node)
            logger.debug(f"Extracted table with {len(rows)} rows")
            
        except Exception as e:
            logger.warning(f"Failed to extract table: {e}")
    
    def parse_slide_metadata(self, prs_slide) -> Dict[str, Any]:
        """Extract metadata about slide layout and structure."""
        return {
            "layout_name": prs_slide.slide_layout.name if prs_slide.slide_layout else "default",
            "has_notes": prs_slide.has_notes_slide,
            "shape_count": len(prs_slide.shapes),
        }
    
    @staticmethod
    def get_presentation_info(pptx_path: str) -> Dict[str, Any]:
        """Get high-level information about presentation."""
        presentation = Presentation(pptx_path)
        
        return {
            "slide_count": len(presentation.slides),
            "dimensions": {
                "width": presentation.slide_width.inches,
                "height": presentation.slide_height.inches,
            },
            "has_master": presentation.slide_master is not None,
        }
